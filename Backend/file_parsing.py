from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

def parse_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def parse_docx(file_path):
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def parse_xlsx(file_path):
    text = ""
    workbook = load_workbook(file_path)
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
    return text

def parse_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
